"""
Knowledge router - handles knowledge source upload, processing, and management.
"""
import os
import uuid
from fastapi import APIRouter, HTTPException, UploadFile, File, Form, BackgroundTasks, Depends
from database import get_supabase
from models.schemas import (
    KnowledgeSourceCreate,
    KnowledgeSourceResponse,
    KnowledgeSourceListResponse,
    ProcessStatusResponse,
)
from services.chunking import split_text, split_documents
from services.embeddings import add_documents_to_vector_store, delete_vectors_by_source
from dependencies.auth import get_current_user_id, require_admin 

router = APIRouter(prefix="/api", tags=["knowledge"])

@router.get("/admin/knowledge")
async def admin_list_all_knowledge(
    user_id: str = Depends(require_admin)  # Requires admin role
):
    """Admin can see all knowledge sources"""
    supabase = get_supabase()
    # No user_id filter - see all users' data
    result = supabase.table("knowledge_sources").select("*").order("created_at", desc=True).execute()
    return result.data


@router.get("/knowledge", response_model=KnowledgeSourceListResponse)
async def list_knowledge_sources(
    search: str = None, 
    type: str = None, 
    status: str = None,
    user_id: str = Depends(get_current_user_id)
):
    """List all knowledge sources for current user."""
    supabase = get_supabase()

    query = supabase.table("knowledge_sources").select("*").eq("user_id", user_id).order("created_at", desc=True)

    if search:
        query = query.ilike("title", f"%{search}%")
    if type:
        query = query.eq("type", type)
    if status:
        query = query.eq("status", status)

    result = query.execute()

    sources = [
        KnowledgeSourceResponse(
            id=s["id"],
            title=s["title"],
            type=s["type"],
            status=s["status"],
            file_path=s.get("file_path"),
            file_size=s.get("file_size"),
            document_count=s.get("document_count", 0),
            created_at=s["created_at"],
            updated_at=s["updated_at"],
        )
        for s in result.data
    ]

    return KnowledgeSourceListResponse(sources=sources, total=len(sources))


@router.get("/knowledge/{source_id}", response_model=KnowledgeSourceResponse)
async def get_knowledge_source(
    source_id: str,
    user_id: str = Depends(get_current_user_id)
):
    """Get a single knowledge source (must belong to current user)."""
    supabase = get_supabase()

    result = supabase.table("knowledge_sources").select("*").eq("id", source_id).eq("user_id", user_id).execute()
    
    if not result.data:
        raise HTTPException(status_code=404, detail="Knowledge source not found")
    
    s = result.data[0]
    return KnowledgeSourceResponse(
        id=s["id"],
        title=s["title"],
        type=s["type"],
        status=s["status"],
        file_path=s.get("file_path"),
        file_size=s.get("file_size"),
        document_count=s.get("document_count", 0),
        created_at=s["created_at"],
        updated_at=s["updated_at"],
    )


@router.post("/knowledge/upload")
async def upload_knowledge_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    title: str = Form(...),
    type: str = Form(default="document"),
    user_id: str = Depends(get_current_user_id)
):
    """
    Upload a file to Supabase Storage and process it for RAG.
    """
    supabase = get_supabase()
    source_id = str(uuid.uuid4())

    # Upload file to Supabase Storage
    file_content = await file.read()
    file_extension = os.path.splitext(file.filename)[1] if file.filename else ""
    storage_path = f"{user_id}/{source_id}{file_extension}"

    try:
        supabase.storage.from_("knowledge-files").upload(storage_path, file_content, {"content-type": file.content_type})
    except Exception as e:
        if "already exists" not in str(e):
            raise HTTPException(status_code=500, detail=f"Storage upload failed: {e}")

    # Create knowledge source record with user_id
    supabase.table("knowledge_sources").insert({
        "id": source_id,
        "title": title or file.filename,
        "type": type,
        "status": "processing",
        "user_id": user_id,
        "file_path": storage_path,
        "file_size": f"{len(file_content) / 1024:.1f} KB",
        "document_count": 0,
    }).execute()

    # Fetch the created record
    source_result = supabase.table("knowledge_sources").select('*').eq('id', source_id).execute()
    source_data = source_result.data[0] if source_result.data else None

    # Schedule background processing
    background_tasks.add_task(process_uploaded_file, source_id, user_id, file_content, file_extension)

    return source_data


@router.post("/knowledge/add-text")
async def add_text_knowledge(
    background_tasks: BackgroundTasks,
    data: KnowledgeSourceCreate,
    user_id: str = Depends(get_current_user_id)
):
    """
    Add text content directly as a knowledge source.
    """
    supabase = get_supabase()
    source_id = str(uuid.uuid4())

    if not data.content:
        raise HTTPException(status_code=400, detail="Content is required for text/faq type")

    # Create knowledge source record with user_id
    supabase.table("knowledge_sources").insert({
        "id": source_id,
        "title": data.title,
        "type": data.type.value,
        "status": "processing",
        "user_id": user_id,
        "document_count": 0,
        "metadata": {"content_length": len(data.content)},
    }).execute()

    # Fetch the created record
    source_result = supabase.table("knowledge_sources").select('*').eq('id', source_id).execute()
    source_data = source_result.data[0] if source_result.data else None

    # Schedule background processing
    background_tasks.add_task(
        process_text_content,
        source_id,
        user_id,
        data.content,
        data.type.value,
    )

    return source_data


@router.delete("/knowledge/{source_id}")
async def delete_knowledge_source(
    source_id: str,
    user_id: str = Depends(get_current_user_id)
):
    """Delete a knowledge source and its vectors (must belong to current user)."""
    supabase = get_supabase()

    # Get source info and verify ownership
    source_result = supabase.table("knowledge_sources").select("file_path").eq("id", source_id).eq("user_id", user_id).execute()
    
    if not source_result.data:
        raise HTTPException(status_code=404, detail="Knowledge source not found")
    
    source = source_result.data[0]

    # Delete from vector store
    delete_vectors_by_source(source_id)

    # Delete file from storage
    if source.get("file_path"):
        try:
            supabase.storage.from_("knowledge-files").remove([source["file_path"]])
        except Exception:
            pass  # File might not exist

    # Delete from database
    supabase.table("knowledge_sources").delete().eq("id", source_id).execute()

    return {"message": "Knowledge source deleted"}


# --- Background Processing Functions ---

def process_uploaded_file(source_id: str, user_id: str, file_content: bytes, file_extension: str):
    """Process an uploaded file: extract text → chunk → embed → store."""
    supabase = get_supabase()

    try:
        # Extract text from file
        text = extract_text_from_file(file_content, file_extension)

        if not text or len(text.strip()) < 10:
            supabase.table("knowledge_sources").update({
                "status": "error",
            }).eq("id", source_id).execute()
            return

        # Chunk the text
        chunks = split_text(text)

        # Store in vector database with user_id in metadata
        documents = [{
            "content": chunk,
            "metadata": {
                "user_id": user_id,
                "source_id": source_id,
            },
        } for chunk in chunks]
        chunks_added = add_documents_to_vector_store(documents, source_id)

        # Update source status
        supabase.table("knowledge_sources").update({
            "status": "active",
            "document_count": chunks_added,
        }).eq("id", source_id).execute()

        print(f"Processed source {source_id}: {chunks_added} chunks created")

    except Exception as e:
        print(f"Error processing source {source_id}: {e}")
        supabase.table("knowledge_sources").update({
            "status": "error",
        }).eq("id", source_id).execute()


def process_text_content(source_id: str, user_id: str, content: str, source_type: str):
    """Process text/FAQ content: chunk → embed → store."""
    supabase = get_supabase()

    try:
        if source_type == "faq":
            # Parse FAQ format
            import re
            qa_pairs = re.split(r'\n(?=(?:Q:|Question:|A:|Answer:))', content)
            documents = []
            for qa in qa_pairs:
                if qa.strip():
                    documents.append({
                        "content": qa.strip(),
                        "metadata": {
                            "user_id": user_id,
                            "source_id": source_id,
                            "type": "faq",
                        },
                    })
        else:
            # Regular text chunking
            chunks = split_text(content)
            documents = [{
                "content": chunk,
                "metadata": {
                    "user_id": user_id,
                    "source_id": source_id,
                },
            } for chunk in chunks]

        # Store in vector database
        chunks_added = add_documents_to_vector_store(documents, source_id)

        # Update source status
        supabase.table("knowledge_sources").update({
            "status": "active",
            "document_count": chunks_added,
        }).eq("id", source_id).execute()

        print(f"Processed text source {source_id}: {chunks_added} chunks created")

    except Exception as e:
        print(f"Error processing text source {source_id}: {e}")
        supabase.table("knowledge_sources").update({
            "status": "error",
        }).eq("id", source_id).execute()


def extract_text_from_file(file_content: bytes, extension: str) -> str:
    """Extract text content from various file types."""
    extension = extension.lower().lstrip(".")

    if extension == "txt":
        import chardet
        detected = chardet.detect(file_content)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        return file_content.decode(encoding, errors="replace")

    elif extension == "pdf":
        import io
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text

    elif extension == "docx":
        import io
        from docx import Document
        doc = Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs if para.text])

    elif extension == "pptx":
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif extension in ["md", "markdown"]:
        return file_content.decode("utf-8", errors="replace")

    elif extension == "csv":
        return file_content.decode("utf-8", errors="replace")

    elif extension == "json":
        import json
        data = json.loads(file_content.decode("utf-8"))
        if isinstance(data, list):
            return "\n".join([json.dumps(item, indent=2) for item in data])
        return json.dumps(data, indent=2)

    else:
        return file_content.decode("utf-8", errors="replace")